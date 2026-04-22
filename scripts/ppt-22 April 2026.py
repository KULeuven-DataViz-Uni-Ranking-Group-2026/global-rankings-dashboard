from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- 1. Define Theme Colors (KU Leuven Navy Blue & Greys) ---
theme_blue = RGBColor(30, 64, 124)   # KU Leuven Navy Blue
text_dark = RGBColor(50, 50, 50)     # Dark Grey for main text
text_light = RGBColor(100, 100, 100) # Lighter Grey for secondary text

# --- 2. Initialize Presentation ---
prs = Presentation()

# Helper function to add the KU Leuven blue strip at the top of a slide
def add_ku_leuven_strip(slide, prs):
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        0, 0, 
        prs.slide_width, Inches(0.15)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = theme_blue
    strip.line.fill.background() # Remove border

# Helper function to style slide titles
def style_title(slide, title_text):
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_frame = title_shape.text_frame
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = theme_blue
            run.font.name = 'Cambria'
            run.font.bold = True

# Helper function to add bullets to a specific shape
def add_bullets_to_shape(shape, bullet_list):
    tf = shape.text_frame
    tf.clear() 
    for i, point in enumerate(bullet_list):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.color.rgb = text_dark
        p.font.name = 'Cambria'
        p.font.size = Pt(20) 

# --- 3. Build the Slides ---

# SLIDE 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
add_ku_leuven_strip(slide1, prs) 
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Prestige vs. Performance:\nVisualizing Biases in Global University Rankings"
title.text_frame.paragraphs[0].runs[0].font.color.rgb = theme_blue
title.text_frame.paragraphs[0].runs[0].font.name = 'Cambria'
title.text_frame.paragraphs[0].runs[0].font.bold = True

subtitle.text = "Data Visualisation [G0R05a] - Intermediate Presentation\nTeam: Renlong Yin, Lei Pei, Victor Kao, Szabó Gergely, Kawtar Darkaoui, Deborah Adelakun"
for p in subtitle.text_frame.paragraphs:
    for run in p.runs:
        run.font.color.rgb = text_light
        run.font.name = 'Cambria'


# SLIDE 2: The Problem & Target Audience
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
add_ku_leuven_strip(slide2, prs) 
style_title(slide2, "The Problem & Target Audience")
bullets_2 = [
    "Target Audience: Prospective international Master's/PhD students.",
    "The Problem: Students often rely blindly on QS Rankings, which heavily favor historical 'prestige' and established Western brands.",
    "The Solution: We compare QS (Reputation) against THE (Research Output) to find the 'Hidden Gems'.",
    "Goal: Help students find universities with world-class research infrastructure that are currently undervalued by general reputation."
]
add_bullets_to_shape(slide2.shapes.placeholders[1], bullets_2)


# SLIDE 3: The Data (No Web Scraping)
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
add_ku_leuven_strip(slide3, prs)
style_title(slide3, "The Data Sources")
bullets_3 = [
    "Times Higher Education (THE) 2016-2026: Provides longitudinal data on research environment, quality, and industry impact.",
    "QS World Rankings 2026: Provides data on Employer and Academic Reputation.",
    "World Bank Open Data API: Live GDP and economic data to analyze the socio-economic drivers of ranking success.",
    "Data Feasibility: We are utilizing robust Kaggle datasets and official APIs rather than scraping, ensuring high data reliability."
]
add_bullets_to_shape(slide3.shapes.placeholders[1], bullets_3)


# SLIDE 4: Proof of Concept / Initial Insights (TWO CONTENT)
slide4 = prs.slides.add_slide(prs.slide_layouts[3])
add_ku_leuven_strip(slide4, prs)
style_title(slide4, "Proof of Concept: Methodology Clash")
bullets_4 = [
    "Our initial Python data exploration confirms our hypothesis.",
    "We identified a massive 'Prestige vs. Performance' gap.",
    "A dense cluster of Asian and European technical universities show elite 'Industry Impact' (THE) but very poor 'Employer Reputation' (QS).",
    "--> See graph to the right."
]
add_bullets_to_shape(slide4.shapes.placeholders[1], bullets_4) 
slide4.shapes.placeholders[2].text = "[Click the picture icon here to insert your Seaborn Scatter Plot]" 


# SLIDE 5: Socio-Economic Drivers (TWO CONTENT - UPDATED)
slide5 = prs.slides.add_slide(prs.slide_layouts[3])
add_ku_leuven_strip(slide5, prs)
style_title(slide5, "Socio-Economic Drivers: GDP vs. Rankings")
bullets_5 = [
    "Data Augmentation: We integrated live World Bank API data to compare a country's wealth against its median university rank.",
    "General Trend: National wealth strongly correlates with elite academic performance (e.g., US, Switzerland, Singapore).",
    "Key Insight: China is a massive outlier. It achieves top-tier rankings despite a significantly lower GDP per capita than its Western rivals.",
    "--> See graph to the right."
]
add_bullets_to_shape(slide5.shapes.placeholders[1], bullets_5) 
slide5.shapes.placeholders[2].text = "[Click the picture icon here to insert your GDP vs Rank Scatter Plot]" 


# SLIDE 6: Questions for Feedback
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
add_ku_leuven_strip(slide6, prs)
style_title(slide6, "Questions for Feedback")
bullets_6 = [
    "1. Visual Clutter: Our dataset has over 800 universities. To avoid violating Gestalt principles in our Scatter Plot, do you recommend semantic zooming or forcing the user to filter by region first?",
    "2. Narrative Flow: Since we pivoted to a student-focused narrative, would a 'Scrollytelling' approach be better than a static dashboard?",
    "3. Encodings: What is the best way to visually encode the 'Gap' between a university's QS rank and THE rank?"
]
add_bullets_to_shape(slide6.shapes.placeholders[1], bullets_6)

# --- 4. Save the Presentation ---
output_filename = "KU_Leuven_DataViz_Presentation.pptx"
prs.save(output_filename)
print(f"Presentation successfully generated and saved as '{output_filename}'")